#!/usr/bin/env python3
"""Office Editor tools for Hermes -- Word/Excel/PPT creation, editing & preview.

Registers tools in the ``office_editor`` toolset that let the Hermes agent:
  - Create new Word (.docx), Excel (.xlsx), PPT (.pptx) documents
  - Open existing Office files from disk
  - Edit document content (text, cells, slides, formatting)
  - Save documents to disk
  - Show documents in the desktop preview pane (real-time WYSIWYG)
  - Query the editor pool status

Architecture:
    Agent tool call
        -> office_editor_tool handler
            -> office_mcp_client.call(tool_name, args)
                -> HTTP POST /mcp  (JSON-RPC 2.0)
                    -> editor_sdk.exe (local binary)
            -> return JSON result

The editor_sdk binary is managed by office_sdk_manager (auto-start on first use).

Tool list (7 high-level tools registered with Hermes):
    1. office_create       -- Create a new doc/sheet/slide
    2. office_open         -- Open an existing file
    3. office_edit         -- Execute any MCP edit operation (unified gateway)
    4. office_save         -- Save a document to disk
    5. office_preview      -- Show a document in the desktop preview pane
    6. office_status       -- Query editor pool status
    7. office_list_tools   -- List all 199 available MCP tools with schemas

Note: each ``registry.register`` below is a top-level statement (not a loop)
so the registry's AST-based tool discovery picks this module up.
"""

import json
import logging
import os
import shlex
import shutil
from typing import Any

from tools.registry import registry, tool_error, tool_result
from tools.office_sdk_manager import sdk_manager
from tools.office_mcp_client import mcp_client
from tools.office_sdk_bridge import resolve_open_file_id, wait_for_editor_ready

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Availability check
# ---------------------------------------------------------------------------
def _sdk_available() -> bool:
    """True when the editor_sdk binary is reachable."""
    from tools.office_sdk_manager import _find_binary
    return _find_binary() is not None


def _officecli_available() -> bool:
    """True when the officecli binary is on PATH (or OFFICE_CLI_COMMAND set)."""
    from tools.office_cli_tool import _OFFICE_CLI_COMMAND
    return shutil.which(_OFFICE_CLI_COMMAND) is not None


def _check_office_engine() -> bool:
    """Availability gate: editor_sdk **or** officecli present.

    The 7 ``office_editor`` tools degrade gracefully: editor_sdk powers the
    rich MCP ops when installed; otherwise the handlers route through
    officecli (path-based create/add/set commands) as the tier-3 fallback.
    """
    return _sdk_available() or _officecli_available()


def _run_officecli(command: str, file_path: str | None = None) -> dict:
    """Run a single officecli command through :mod:`tools.office_cli_tool`."""
    from tools.office_cli_tool import run_office_cli_command
    return run_office_cli_command(command, file_path=file_path)

# ---------------------------------------------------------------------------
# Helper: determine doc_type from file extension
# ---------------------------------------------------------------------------
_EXT_MAP = {
    ".doc": "doc", ".docx": "doc", ".dot": "doc", ".wps": "doc", ".wpt": "doc", ".docm": "doc",
    ".xls": "sheet", ".xlsx": "sheet", ".xlt": "sheet", ".csv": "sheet", ".tsv": "sheet", ".xlsm": "sheet",
    ".ppt": "slide", ".pptx": "slide", ".pps": "slide", ".pot": "slide", ".pptm": "slide",
}

def _doc_type_from_path(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    return _EXT_MAP.get(ext, "doc")

def _create_tool_name(doc_type: str) -> str:
    return {"doc": "create_doc", "sheet": "create_sheet", "slide": "create_slide"}.get(doc_type, "create_doc")

def _wait_for_editor_ready(file_id: str, doc_type: str = "doc", timeout: float = 15.0) -> bool:
    """Wait until editor_sdk can actually serve an edit on this editor.

    Thin wrapper over :func:`tools.office_sdk_bridge.wait_for_editor_ready`
    passing this module's ``mcp_client`` singleton.
    """
    return wait_for_editor_ready(mcp_client, file_id, doc_type, timeout)


def _resolve_open_file_id(file_path: str, timeout: float = 10.0) -> str:
    """Find the file_id editor_sdk assigned to an ``open_file`` editor.

    Thin wrapper over :func:`tools.office_sdk_bridge.resolve_open_file_id`
    passing this module's ``sdk_manager`` singleton.
    """
    return resolve_open_file_id(sdk_manager, file_path, timeout)

# ---------------------------------------------------------------------------
# Tool 1: office_create
# ---------------------------------------------------------------------------
OFFICE_CREATE_SCHEMA = {
    "name": "office_create",
    "description": (
        "Create a new Word (.docx), Excel (.xlsx), or PowerPoint (.pptx) document "
        "from a blank template. Returns a file_id that can be used with "
        "office_edit, office_save, and office_preview.\n\n"
        "Examples:\n"
        "  office_create(doc_type=\"doc\", file_path=\"/tmp/report.docx\")\n"
        "  office_create(doc_type=\"sheet\", file_path=\"/tmp/data.xlsx\")\n"
        "  office_create(doc_type=\"slide\", file_path=\"/tmp/slides.pptx\")"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "doc_type": {
                "type": "string",
                "enum": ["doc", "sheet", "slide"],
                "description": "Document type: doc (Word), sheet (Excel), or slide (PowerPoint).",
            },
            "file_path": {
                "type": "string",
                "description": "Absolute path where the document will be saved (used for naming; actual save happens via office_save).",
            },
        },
        "required": ["doc_type", "file_path"],
    },
}

def _handle_office_create(args: dict, **kwargs) -> str:
    doc_type = args.get("doc_type", "doc")
    file_path = args.get("file_path", "")
    if not file_path:
        return tool_error("file_path is required")

    if _sdk_available():
        try:
            result = mcp_client.call(_create_tool_name(doc_type), {"file_path": file_path})
        except RuntimeError as e:
            return tool_error(str(e))

        file_id = result.get("file_id", "")
        if not file_id:
            return tool_error("No file_id returned from editor_sdk")

        # The editor registers asynchronously; wait until it can actually serve an
        # edit before returning so a subsequent office_edit does not hit
        # "document is not open" / "No workbook open".
        _wait_for_editor_ready(file_id, doc_type)

        # Try to open in the desktop preview pane
        _try_open_preview(file_id, doc_type)

        return tool_result(
            success=True,
            file_id=file_id,
            doc_type=doc_type,
            file_path=result.get("file_path", file_path),
            engine="editor_sdk",
            message=f"Created {doc_type} document. Use office_edit to modify content, office_save to save to disk.",
        )

    # officecli fallback: create a blank file on disk; file_id is the path.
    path = os.path.abspath(file_path)
    result = _run_officecli(f"officecli create {shlex.quote(path)}", file_path=path)
    if not result.get("success"):
        return tool_error(result.get("stderr") or result.get("error") or "officecli create failed")
    return tool_result(
        success=True,
        file_id=path,
        doc_type=doc_type,
        file_path=path,
        engine="officecli",
        message=f"Created {doc_type} document via officecli. Use office_cli_command to modify content, office_save to flush to disk.",
    )

# ---------------------------------------------------------------------------
# Tool 2: office_open
# ---------------------------------------------------------------------------
OFFICE_OPEN_SCHEMA = {
    "name": "office_open",
    "description": (
        "Open an existing Office file (.docx/.xlsx/.pptx and variants) from disk "
        "for editing. Returns a file_id for subsequent edit/save/preview operations.\n\n"
        "Example:\n"
        "  office_open(file_path=\"/home/user/report.docx\")"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Absolute path to the file on disk.",
            },
        },
        "required": ["file_path"],
    },
}

def _handle_office_open(args: dict, **kwargs) -> str:
    file_path = args.get("file_path", "")
    if not file_path:
        return tool_error("file_path is required")

    # Normalize path
    file_path = os.path.abspath(file_path)
    if not os.path.exists(file_path):
        return tool_error(f"File not found: {file_path}")

    doc_type = _doc_type_from_path(file_path)

    if _sdk_available():
        try:
            result = mcp_client.call("open_file", {"file_path": file_path, "file_type": doc_type})
        except RuntimeError as e:
            return tool_error(str(e))

        # open_file is streaming and returns no file_id on the envelope; the editor
        # registers in the pool asynchronously, keyed by its file path.
        file_id = result.get("file_id") or _resolve_open_file_id(file_path)
        if not file_id:
            return tool_error("No file_id returned from editor_sdk")

        _wait_for_editor_ready(file_id, doc_type)
        _try_open_preview(file_id, doc_type)

        return tool_result(
            success=True,
            file_id=file_id,
            doc_type=doc_type,
            file_path=file_path,
            engine="editor_sdk",
            message=f"Opened {doc_type} document. Use office_edit to modify content.",
        )

    # officecli fallback: officecli operates on the path directly.
    return tool_result(
        success=True,
        file_id=file_path,
        doc_type=doc_type,
        file_path=file_path,
        engine="officecli",
        message=f"Opened {doc_type} document (officecli). Use office_cli_command to modify content.",
    )

# ---------------------------------------------------------------------------
# Tool 3: office_edit (unified gateway to all 199 MCP tools)
# ---------------------------------------------------------------------------
OFFICE_EDIT_SCHEMA = {
    "name": "office_edit",
    "description": (
        "Execute any editor_sdk MCP edit operation on an open document. "
        "This is the unified gateway to all 199 document editing tools.\n\n"
        "Common operations:\n"
        "  - doc_insert_text: Insert text at a position in a Word doc\n"
        "  - doc_replace_range: Replace a text range in a Word doc\n"
        "  - doc_set_font: Set font properties for a range\n"
        "  - doc_insert_paragraph: Insert a new paragraph\n"
        "  - doc_get_outline: Get the document outline/headings\n"
        "  - sheet_set_cell_value: Set a cell value in Excel\n"
        "  - sheet_set_range_value: Set a range of cells\n"
        "  - sheet_get_sheet_info: Get sheet names and IDs\n"
        "  - sheet_set_column_width: Set column width\n"
        "  - slide_add_slide: Add a new slide\n"
        "  - slide_set_text: Set text in a slide placeholder\n"
        "  - slide_add_textbox: Add a text box\n"
        "  - slide_set_background: Set slide background\n\n"
        "Use office_list_tools to see all available operations with their schemas.\n\n"
        "Example:\n"
        "  office_edit(\n"
        "    file_id=\"new_doc_xxx\",\n"
        "    operation=\"doc_insert_text\",\n"
        "    arguments={\"idx\": 0, \"text\": \"Hello World\"}\n"
        "  )"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "file_id": {
                "type": "string",
                "description": "The file_id from office_create or office_open.",
            },
            "operation": {
                "type": "string",
                "description": (
                    "The MCP tool name to execute. "
                    "Examples: doc_insert_text, sheet_set_cell_value, slide_add_slide. "
                    "Use office_list_tools to see all 199 operations."
                ),
            },
            "arguments": {
                "type": "object",
                "description": "Arguments for the specific operation. The file_id is auto-injected.",
                "properties": {},
                "additionalProperties": True,
            },
        },
        "required": ["file_id", "operation"],
    },
}

def _handle_office_edit(args: dict, **kwargs) -> str:
    file_id = args.get("file_id", "")
    operation = args.get("operation", "")
    op_args = args.get("arguments", {})

    if not file_id:
        return tool_error("file_id is required")
    if not operation:
        return tool_error("operation is required (e.g. 'doc_insert_text')")

    if _sdk_available():
        # Inject file_id into the operation arguments
        op_args = dict(op_args)
        op_args["file_id"] = file_id

        try:
            result = mcp_client.call(operation, op_args)
        except RuntimeError as e:
            return tool_error(str(e))

        # Extract the text content for the result
        content = result.get("content", [])
        text_content = ""
        if content and isinstance(content, list):
            text_content = content[0].get("text", "")

        # Try to parse as JSON for structured results
        parsed = None
        if text_content:
            try:
                parsed = json.loads(text_content)
            except json.JSONDecodeError:
                parsed = None

        return tool_result(
            success=True,
            operation=operation,
            file_id=file_id,
            engine="editor_sdk",
            result=parsed if parsed is not None else text_content,
        )

    # officecli fallback: the 199 SDK MCP ops do not map 1:1 onto officecli's
    # path-based `add`/`set --type/--prop` CLI, so we do not fabricate a
    # translation. Direct the model to the office_cli_command tool, which
    # speaks officecli's native syntax and is available in officecli mode.
    return tool_result(
        success=False,
        operation=operation,
        file_id=file_id,
        engine="officecli",
        message=(
            "officecli fallback mode: this operation is not mapped. "
            "Use the office_cli_command tool with officecli's native syntax "
            "to edit the file (e.g. `officecli add <file> /body --type paragraph "
            "--prop text=\"...\"`)."
        ),
    )

# ---------------------------------------------------------------------------
# Tool 4: office_save
# ---------------------------------------------------------------------------
OFFICE_SAVE_SCHEMA = {
    "name": "office_save",
    "description": (
        "Save a document to disk. The file is saved to the path specified "
        "during creation/opening, or to a custom path via save_path.\n\n"
        "Example:\n"
        "  office_save(file_id=\"new_doc_xxx\", save_path=\"/tmp/output.docx\")"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "file_id": {
                "type": "string",
                "description": "The file_id from office_create or office_open.",
            },
            "save_path": {
                "type": "string",
                "description": "Optional custom save path. If omitted, saves to the original path.",
            },
        },
        "required": ["file_id"],
    },
}

def _handle_office_save(args: dict, **kwargs) -> str:
    file_id = args.get("file_id", "")
    save_path = args.get("save_path", "")

    if not file_id:
        return tool_error("file_id is required")

    if _sdk_available():
        # NOTE: editor_sdk's save_file tool expects the destination under the
        # "file_path" argument (not "save_path").
        call_args = {"file_id": file_id}
        if save_path:
            call_args["file_path"] = os.path.abspath(save_path)

        try:
            text = mcp_client.call_text("save_file", call_args)
        except RuntimeError as e:
            return tool_error(str(e))

        return tool_result(
            success=True,
            file_id=file_id,
            engine="editor_sdk",
            message=text,
        )

    # officecli fallback: flush the file to disk. In officecli mode file_id is
    # the on-disk path (from office_open/office_create).
    path = os.path.abspath(save_path or file_id)
    result = _run_officecli(f"officecli save {shlex.quote(path)}", file_path=path)
    if not result.get("success"):
        return tool_error(result.get("stderr") or result.get("error") or "officecli save failed")
    return tool_result(
        success=True,
        file_id=file_id,
        engine="officecli",
        message="Saved via officecli (file flushed to disk).",
    )

# ---------------------------------------------------------------------------
# Tool 5: office_preview
# ---------------------------------------------------------------------------
OFFICE_PREVIEW_SCHEMA = {
    "name": "office_preview",
    "description": (
        "Show a document in the Hermes desktop preview pane for real-time "
        "WYSIWYG editing. The user sees all changes made via office_edit "
        "in real time.\n\n"
        "Example:\n"
        "  office_preview(file_id=\"new_doc_xxx\", doc_type=\"doc\")"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "file_id": {
                "type": "string",
                "description": "The file_id from office_create or office_open.",
            },
            "doc_type": {
                "type": "string",
                "enum": ["doc", "sheet", "slide"],
                "description": "Document type for the preview renderer.",
            },
        },
        "required": ["file_id", "doc_type"],
    },
}

def _handle_office_preview(args: dict, **kwargs) -> str:
    file_id = args.get("file_id", "")
    doc_type = args.get("doc_type", "doc")

    if not file_id:
        return tool_error("file_id is required")

    if _sdk_available():
        # editor_sdk mode is read-only for the human: do NOT pass editable /
        # file_path, so sheets route to the SDK read-only cloud view instead of
        # the Univer WYSIWYG editor.
        url = sdk_manager.get_preview_url(file_id, doc_type)

        if not _try_open_preview(file_id, doc_type):
            # Not in desktop mode -- return the URL so the user can open it manually
            return tool_result(
                success=True,
                file_id=file_id,
                preview_url=url,
                engine="editor_sdk",
                message=(
                    "Preview URL generated. Open it in a browser to view the document. "
                    f"URL: {url}"
                ),
            )

        return tool_result(
            success=True,
            file_id=file_id,
            preview_url=url,
            engine="editor_sdk",
            message="Document opened in the preview pane. Edits via office_edit will appear in real time.",
        )

    # officecli fallback: start the watch server and return its URL.
    from tools.office_cli_tool import start_office_preview
    result = start_office_preview(file_id)
    if "url" in result:
        return tool_result(
            success=True,
            file_id=file_id,
            preview_url=result["url"],
            engine="officecli",
            message=f"Preview opened via officecli. URL: {result['url']}",
        )
    return tool_error(result.get("message", "officecli preview failed"))

# ---------------------------------------------------------------------------
# Tool 6: office_status
# ---------------------------------------------------------------------------
OFFICE_STATUS_SCHEMA = {
    "name": "office_status",
    "description": (
        "Query the editor_sdk status: running port, health, and list of "
        "currently open document editors with their file_ids and save state."
    ),
    "parameters": {
        "type": "object",
        "properties": {},
    },
}

def _handle_office_status(args: dict, **kwargs) -> str:
    if _sdk_available():
        try:
            status = sdk_manager.get_editor_status()
            return tool_result(
                success=True,
                port=sdk_manager.port,
                healthy=sdk_manager.health_check(),
                engine="editor_sdk",
                **status,
            )
        except Exception as e:
            return tool_error(str(e))

    # officecli fallback: report the watch-server sessions.
    from tools.office_cli_tool import _sessions, _sessions_lock
    with _sessions_lock:
        sessions = [
            {"file_path": fp, "running": s["process"].poll() is None}
            for fp, s in _sessions.items()
        ]
    return tool_result(
        success=True,
        engine="officecli",
        sessions=sessions,
        open_count=len(sessions),
    )

# ---------------------------------------------------------------------------
# Tool 7: office_list_tools
# ---------------------------------------------------------------------------
OFFICE_LIST_TOOLS_SCHEMA = {
    "name": "office_list_tools",
    "description": (
        "List all available editor_sdk MCP tools (199 total) with their "
        "input schemas. Use this to discover the exact parameter names "
        "and types for the 'operation' and 'arguments' fields of office_edit.\n\n"
        "Optionally filter by prefix: 'doc', 'sheet', 'slide', or omit to list all."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "prefix": {
                "type": "string",
                "description": "Filter by tool name prefix: 'doc', 'sheet', 'slide', 'create', 'save', 'open', etc. Omit to list all.",
            },
        },
    },
}

def _handle_office_list_tools(args: dict, **kwargs) -> str:
    prefix = args.get("prefix", "")

    if _sdk_available():
        try:
            tools = mcp_client.list_tools()
        except RuntimeError as e:
            return tool_error(str(e))

        if prefix:
            tools = [t for t in tools if t["name"].startswith(prefix)]

        # Return a compact summary
        summary = []
        for t in tools:
            schema = t.get("inputSchema", {})
            required = schema.get("required", [])
            properties = schema.get("properties", {})
            props_str = ", ".join(
                f"{k}:{v.get('type', 'any')}" for k, v in properties.items()
                if k != "file_id"  # file_id is auto-injected
            )
            summary.append({
                "name": t["name"],
                "title": t.get("title", t.get("annotations", {}).get("title", "")),
                "description": t.get("description", "")[:120],
                "required": required,
                "properties": props_str,
            })

        return tool_result(
            success=True,
            engine="editor_sdk",
            total=len(tools),
            tools=summary,
        )

    # officecli fallback: no MCP schema to enumerate.
    return tool_result(
        success=True,
        engine="officecli",
        total=0,
        tools=[],
        message=(
            "editor_sdk unavailable. officecli mode: use the office_cli_command "
            "tool with officecli's native `create`/`add`/`set`/`get` commands "
            "(see the officecli skill)."
        ),
    )

# ---------------------------------------------------------------------------
# Tool 8: office_render
# ---------------------------------------------------------------------------
OFFICE_RENDER_SCHEMA = {
    "name": "office_render",
    "description": (
        "Render an open document to PNG images (and optionally PDF) for visual "
        "quality assurance. Converts via OnlyOffice (Microsoft Office COM "
        "fallback) and rasterizes with pymupdf, revealing layout defects "
        "invisible through text extraction.\n\n"
        "Supported document types: docx, xlsx, pptx (and variants).\n\n"
        "Examples:\n"
        "  office_render(file_id=\"new_doc_xxx\", doc_type=\"doc\")\n"
        "  office_render(file_id=\"new_doc_xxx\", doc_type=\"doc\", format=\"pdf\")\n"
        "  office_render(file_id=\"new_doc_xxx\", doc_type=\"slide\", dpi=200)"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "file_id": {
                "type": "string",
                "description": "The file_id from office_create or office_open.",
            },
            "doc_type": {
                "type": "string",
                "enum": ["doc", "sheet", "slide"],
                "description": "Document type: doc (Word), sheet (Excel), or slide (PowerPoint).",
            },
            "format": {
                "type": "string",
                "enum": ["png", "pdf"],
                "default": "png",
                "description": "Output format: png (page images) or pdf (single file). Default: png.",
            },
            "dpi": {
                "type": "integer",
                "default": 150,
                "description": "Render resolution in DPI. Default: 150.",
            },
            "seal_text": {
                "type": "string",
                "description": (
                    "Optional: audit organization name to stamp a default seal "
                    "onto the first page (cover) of the PDF output. Only applies "
                    "when format='pdf'. Omit to skip stamping."
                ),
            },
        },
        "required": ["file_id", "doc_type"],
    },
}

def _pdf_to_png(pdf_path: str, output_dir: str, dpi: int) -> list[str]:
    """Rasterize PDF pages to PNG via pymupdf (pure pip, no poppler)."""
    import pymupdf

    images: list[str] = []
    with pymupdf.open(pdf_path) as doc:
        for index, page in enumerate(doc):
            img_path = os.path.join(output_dir, f"page-{index + 1:03d}.png")
            page.get_pixmap(dpi=dpi).save(img_path)
            images.append(os.path.abspath(img_path))
    return images


def _handle_office_render(args: dict, **kwargs) -> str:
    file_id = args.get("file_id", "")
    doc_type = args.get("doc_type", "doc")
    fmt = args.get("format", "png")
    dpi = args.get("dpi", 150)
    seal_text = (args.get("seal_text") or "").strip() or None

    if not file_id:
        return tool_error("file_id is required")

    # Resolve the file path from the editor pool
    file_path = None
    if _sdk_available():
        try:
            status = sdk_manager.get_editor_status()
            for ed in status.get("open_editors", []):
                if ed.get("file_id") == file_id:
                    file_path = ed.get("file_path") or ed.get("save_path")
                    break
        except Exception:
            pass

    if not file_path and _sdk_available():
        # Try saving first to ensure the file exists on disk
        try:
            result = mcp_client.call("save_file", {"file_id": file_id})
            text = (result.get("content") or [{}])[0].get("text", "")
            import json as _json
            data = _json.loads(text) if text else {}
            file_path = data.get("file_path") or data.get("path")
        except Exception:
            pass

    if not file_path:
        return tool_error(
            "Cannot resolve file path for rendering. Save the document first "
            "with office_save, then retry."
        )

    import os
    import tempfile

    if not os.path.exists(file_path):
        return tool_error(f"File not found on disk: {file_path}")

    output_dir = tempfile.mkdtemp(prefix="hermes_render_")

    try:
        # Office → PDF via OnlyOffice ConvertService, COM automation fallback.
        from tools.office_pdf_convert import office_to_pdf
        pdf_path = office_to_pdf(
            file_path,
            doc_type,
            seal_text=seal_text if fmt == "pdf" else None,
        )
    except Exception as exc:
        return tool_error(f"Rendering failed: {exc}")

    try:
        if fmt == "pdf":
            import shutil as _shutil
            final_path = os.path.join(tempfile.gettempdir(), f"hermes_render_{file_id}.pdf")
            _shutil.copy(pdf_path, final_path)
            return tool_result(
                success=True,
                file_id=file_id,
                file_path=final_path,
                format="pdf",
                message=f"Rendered to PDF: {final_path}. Use vision_analyze to inspect.",
            )

        # Convert PDF pages to PNG images via pymupdf.
        images = _pdf_to_png(pdf_path, output_dir, dpi)
        if not images:
            return tool_error("pymupdf produced no output images")

        return tool_result(
            success=True,
            file_id=file_id,
            page_count=len(images),
            dpi=dpi,
            images=images,
            message=(
                f"Rendered {len(images)} page(s) at {dpi} DPI. "
                "Inspect each image with vision_analyze for layout defects."
            ),
        )
    except Exception as e:
        return tool_error(str(e))


# ---------------------------------------------------------------------------
# Tool 9: office_audit
# ---------------------------------------------------------------------------
OFFICE_AUDIT_SCHEMA = {
    "name": "office_audit",
    "description": (
        "Run quality audits on an open document. Checks heading hierarchy, "
        "style consistency, image accessibility, table geometry, and more.\n\n"
        "Audit types:\n"
        "  - headings: Check heading levels, nesting, and numbering consistency\n"
        "  - styles: Check for inconsistent or missing styles\n"
        "  - a11y: Accessibility audit (alt text, table headers, reading order)\n"
        "  - tables: Table geometry audit (widths, indents, cell margins)\n"
        "  - all: Run all available audits\n\n"
        "Examples:\n"
        "  office_audit(file_id=\"new_doc_xxx\", audit_type=\"headings\")\n"
        "  office_audit(file_id=\"new_doc_xxx\", audit_type=\"all\")"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "file_id": {
                "type": "string",
                "description": "The file_id from office_create or office_open.",
            },
            "audit_type": {
                "type": "string",
                "enum": ["headings", "styles", "a11y", "tables", "all"],
                "default": "all",
                "description": "Which audit to run. Default: all.",
            },
        },
        "required": ["file_id"],
    },
}

def _handle_office_audit(args: dict, **kwargs) -> str:
    file_id = args.get("file_id", "")
    audit_type = args.get("audit_type", "all")

    if not file_id:
        return tool_error("file_id is required")

    # Resolve file path
    file_path = None
    if _sdk_available():
        try:
            status = sdk_manager.get_editor_status()
            for ed in status.get("open_editors", []):
                if ed.get("file_id") == file_id:
                    file_path = ed.get("file_path") or ed.get("save_path")
                    break
        except Exception:
            pass

    if not file_path or not os.path.exists(file_path):
        return tool_error("Cannot resolve file path. Save the document first with office_save.")

    import subprocess
    import os as _os

    skill_scripts = _os.path.join(
        _os.path.dirname(_os.path.dirname(_os.path.abspath(__file__))),
        "skills", "productivity", "docx", "scripts"
    )

    results = {}

    def _run_script(script_name, *extra_args):
        script_path = _os.path.join(skill_scripts, script_name)
        if not _os.path.exists(script_path):
            return {"error": f"Script not found: {script_name}"}
        try:
            proc = subprocess.run(
                ["python", script_path, file_path] + list(extra_args),
                capture_output=True, text=True, timeout=30,
            )
            return {"ok": proc.returncode == 0, "stdout": proc.stdout[:2000], "stderr": proc.stderr[:500]}
        except Exception as exc:
            return {"error": str(exc)}

    if audit_type in ("headings", "all"):
        results["headings"] = _run_script("heading_audit.py")
    if audit_type in ("styles", "all"):
        results["styles"] = _run_script("style_lint.py")
    if audit_type in ("a11y", "all"):
        results["a11y"] = _run_script("a11y_audit.py")
    if audit_type in ("tables", "all"):
        results["tables"] = _run_script("table_geometry.py")

    issues_found = sum(
        1 for r in results.values()
        if isinstance(r, dict) and r.get("ok") is False
    )

    return tool_result(
        success=True,
        file_id=file_id,
        audit_type=audit_type,
        issues_found=issues_found,
        results=results,
        message=(
            f"Audit complete. {issues_found} issue(s) found across "
            f"{len(results)} check(s)."
        ),
    )


# ---------------------------------------------------------------------------
# Tool 10: office_template_extract
# ---------------------------------------------------------------------------
OFFICE_TEMPLATE_EXTRACT_SCHEMA = {
    "name": "office_template_extract",
    "description": (
        "Extract design rules from an existing document for use as a template. "
        "Analyzes the document's fonts, colors, heading styles, spacing patterns, "
        "and table formatting to produce a reusable design token map.\n\n"
        "Use this before creating a new document that should match an existing "
        "template's style. The extracted rules can be passed as design guidance "
        "to office_create.\n\n"
        "Examples:\n"
        "  office_template_extract(file_id=\"template_xxx\", format=\"yaml\")\n"
        "  office_template_extract(file_path=\"/path/to/template.docx\")"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "file_id": {
                "type": "string",
                "description": "The file_id from office_open of the template document.",
            },
            "file_path": {
                "type": "string",
                "description": "Alternative: absolute path to a template file on disk. Use this if the template is not yet open.",
            },
        },
    },
}

def _handle_office_template_extract(args: dict, **kwargs) -> str:
    file_id = args.get("file_id", "")
    file_path = args.get("file_path", "")

    if not file_id and not file_path:
        return tool_error("Either file_id or file_path is required")

    if file_id and not file_path:
        # Resolve from pool
        if _sdk_available():
            try:
                status = sdk_manager.get_editor_status()
                for ed in status.get("open_editors", []):
                    if ed.get("file_id") == file_id:
                        file_path = ed.get("file_path") or ed.get("save_path")
                        break
            except Exception:
                pass

    if not file_path or not os.path.exists(file_path):
        return tool_error("Cannot resolve template file path.")

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in (".docx", ".dotx", ".pptx", ".potx"):
        return tool_error(f"Unsupported template format: {ext}. Supported: .docx, .dotx, .pptx, .potx")

    # Use python-docx or python-pptx to extract design tokens
    tokens = {
        "source": os.path.abspath(file_path),
        "format": ext,
    }

    try:
        if ext in (".docx", ".dotx"):
            from docx import Document as DocxDoc
            doc = DocxDoc(file_path)

            # Extract fonts from styles
            fonts = set()
            heading_info = {}
            for style in doc.styles:
                if style.font and style.font.name:
                    fonts.add(style.font.name)
                if style.name and style.name.startswith("Heading"):
                    try:
                        level = int(style.name.split()[-1])
                        heading_info[f"h{level}"] = {
                            "font": style.font.name if style.font else "unknown",
                            "size": str(style.font.size) if style.font and style.font.size else "unknown",
                            "color": str(style.font.color.rgb) if style.font and style.font.color and style.font.color.rgb else "unknown",
                        }
                    except ValueError:
                        pass

            tokens["fonts"] = sorted(fonts)
            tokens["headings"] = heading_info

            # Extract page setup
            for section in doc.sections:
                tokens["page"] = {
                    "width": section.page_width,
                    "height": section.page_height,
                    "margin_top": section.top_margin,
                    "margin_bottom": section.bottom_margin,
                    "margin_left": section.left_margin,
                    "margin_right": section.right_margin,
                }
                break  # first section only

        elif ext in (".pptx", ".potx"):
            from pptx import Presentation as PptxPres
            prs = PptxPres(file_path)

            tokens["slide_count"] = len(prs.slides)
            tokens["slide_width"] = prs.slide_width
            tokens["slide_height"] = prs.slide_height

            # Extract slide layouts
            layouts = []
            for layout in prs.slide_layouts:
                layouts.append({
                    "name": layout.name,
                    "placeholder_count": len(layout.placeholders),
                })
            tokens["slide_layouts"] = layouts

        return tool_result(
            success=True,
            file_path=os.path.abspath(file_path),
            format=ext,
            tokens=tokens,
            message=(
                f"Extracted design tokens from {os.path.basename(file_path)}. "
                "Use these tokens to guide document creation with consistent styling."
            ),
        )
    except ImportError as e:
        return tool_error(f"Missing Python package: {e}. Install python-docx and/or python-pptx.")
    except Exception as e:
        return tool_error(f"Template extraction failed: {e}")


# ---------------------------------------------------------------------------
# Helper: open preview in desktop UI
# ---------------------------------------------------------------------------
def _try_open_preview(file_id: str, doc_type: str) -> bool:
    """Try to open the document in the Hermes desktop preview pane.

    Returns True if the desktop UI accepted the request, False otherwise
    (e.g. running in CLI mode without a desktop GUI).
    """
    try:
        from tools import desktop_ui
        if not desktop_ui.available():
            return False
        url = sdk_manager.get_preview_url(file_id, doc_type, editable=True)
        return desktop_ui.emit("preview.open", {
            "url": url,
            "label": f"Office Editor ({doc_type})",
        })
    except Exception as e:
        logger.debug("Could not open desktop preview: %s", e)
        return False

# ---------------------------------------------------------------------------
# Registration
# ---------------------------------------------------------------------------
# NOTE: each registration is a top-level statement so the registry's AST scan
# (tools.registry._module_registers_tools) detects this module and imports it.
registry.register(
    name="office_create",
    toolset="office_editor",
    schema=OFFICE_CREATE_SCHEMA,
    handler=_handle_office_create,
    check_fn=_check_office_engine,
    requires_env=[],
    is_async=False,
    description=OFFICE_CREATE_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f4c4",
)

registry.register(
    name="office_open",
    toolset="office_editor",
    schema=OFFICE_OPEN_SCHEMA,
    handler=_handle_office_open,
    check_fn=_check_office_engine,
    requires_env=[],
    is_async=False,
    description=OFFICE_OPEN_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f4c4",
)

registry.register(
    name="office_edit",
    toolset="office_editor",
    schema=OFFICE_EDIT_SCHEMA,
    handler=_handle_office_edit,
    check_fn=_check_office_engine,
    requires_env=[],
    is_async=False,
    description=OFFICE_EDIT_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f4c4",
)

registry.register(
    name="office_save",
    toolset="office_editor",
    schema=OFFICE_SAVE_SCHEMA,
    handler=_handle_office_save,
    check_fn=_check_office_engine,
    requires_env=[],
    is_async=False,
    description=OFFICE_SAVE_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f4c4",
)

registry.register(
    name="office_preview",
    toolset="office_editor",
    schema=OFFICE_PREVIEW_SCHEMA,
    handler=_handle_office_preview,
    check_fn=_check_office_engine,
    requires_env=[],
    is_async=False,
    description=OFFICE_PREVIEW_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f4c4",
)

registry.register(
    name="office_status",
    toolset="office_editor",
    schema=OFFICE_STATUS_SCHEMA,
    handler=_handle_office_status,
    check_fn=_check_office_engine,
    requires_env=[],
    is_async=False,
    description=OFFICE_STATUS_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f4c4",
)

registry.register(
    name="office_list_tools",
    toolset="office_editor",
    schema=OFFICE_LIST_TOOLS_SCHEMA,
    handler=_handle_office_list_tools,
    check_fn=_check_office_engine,
    requires_env=[],
    is_async=False,
    description=OFFICE_LIST_TOOLS_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f4c4",
)

registry.register(
    name="office_render",
    toolset="office_editor",
    schema=OFFICE_RENDER_SCHEMA,
    handler=_handle_office_render,
    check_fn=None,  # OnlyOffice/COM/pymupdf resolved at runtime; works without editor_sdk
    requires_env=[],
    is_async=False,
    description=OFFICE_RENDER_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f3de",
)

registry.register(
    name="office_audit",
    toolset="office_editor",
    schema=OFFICE_AUDIT_SCHEMA,
    handler=_handle_office_audit,
    check_fn=None,  # checks at runtime
    requires_env=[],
    is_async=False,
    description=OFFICE_AUDIT_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f50d",
)

registry.register(
    name="office_template_extract",
    toolset="office_editor",
    schema=OFFICE_TEMPLATE_EXTRACT_SCHEMA,
    handler=_handle_office_template_extract,
    check_fn=None,  # checks at runtime
    requires_env=[],
    is_async=False,
    description=OFFICE_TEMPLATE_EXTRACT_SCHEMA.get("description", "").split("\n")[0][:80],
    emoji="\U0001f3a8",
)
