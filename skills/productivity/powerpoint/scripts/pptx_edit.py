#!/usr/bin/env python3
"""Edit a .pptx in place (or save to --output).

Operations (repeatable / combinable):
  --replace-text OLD NEW   Replace text everywhere (slides, tables, notes).
                           Run-level replace preserves formatting when the
                           match sits inside a single run; if PowerPoint has
                           split the match across runs, the whole paragraph
                           is rewritten using the first run's font.
  --chart-data SPEC.json   Replace data of a chart. Spec:
                           {"slide": 0, "chart": 0,
                            "categories": ["Q1", "Q2"],
                            "series": {"North": [1, 2], "South": [3, 4]}}
  --swap-image SLIDE SHAPE_NAME NEW_IMAGE
                           Replace a picture's bits, keeping position/size.
  --remove-slide N         Delete slide at index N (0-based).
  --move-slide FROM TO     Reorder: move slide FROM to position TO
                           (0-based; XML-level _sldIdLst manipulation).
"""
import argparse
import copy
import json
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE


def replace_in_text_frame(text_frame, old, new):
    count = 0
    for para in text_frame.paragraphs:
        joined = "".join(run.text for run in para.runs)
        if old not in joined:
            continue
        if any(old in run.text for run in para.runs):
            # Run-level replace: preserves each run's formatting exactly.
            for run in para.runs:
                if old in run.text:
                    count += run.text.count(old)
                    run.text = run.text.replace(old, new)
        else:
            # Match split across runs -> rewrite paragraph text,
            # keeping only the first run's formatting (documented caveat).
            count += joined.count(old)
            first = para.runs[0]
            first.text = joined.replace(old, new)
            for run in para.runs[1:]:
                run._r.getparent().remove(run._r)
    return count


def iter_text_frames(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
    if slide.has_notes_slide:
        yield slide.notes_slide.notes_text_frame


def replace_text(prs, old, new):
    total = 0
    for slide in prs.slides:
        for tf in iter_text_frames(slide):
            total += replace_in_text_frame(tf, old, new)
    return total


def update_chart(prs, spec_path):
    with open(spec_path, encoding="utf-8") as fh:
        spec = json.load(fh)
    slide = prs.slides[spec.get("slide", 0)]
    charts = [s.chart for s in slide.shapes if s.has_chart]
    if not charts:
        raise SystemExit(f"no chart on slide {spec.get('slide', 0)}")
    chart = charts[spec.get("chart", 0)]
    data = CategoryChartData()
    data.categories = spec["categories"]
    for name, values in spec["series"].items():
        data.add_series(name, values)
    chart.replace_data(data)


def swap_image(prs, slide_idx, shape_name, new_path):
    slide = prs.slides[int(slide_idx)]
    for shape in slide.shapes:
        if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                and shape.name == shape_name):
            image_part, rid = slide.part.get_or_add_image_part(new_path)
            blip = shape._element.blipFill.blip
            blip.set("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed", rid)
            return True
    raise SystemExit(f"no picture named {shape_name!r} on slide {slide_idx}")


def remove_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    slide_id = list(sldIdLst)[int(index)]
    rid = slide_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rid)
    sldIdLst.remove(slide_id)


def move_slide(prs, src, dst):
    """Reorder by moving the <p:sldId> element inside <p:sldIdLst>."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    element = ids[int(src)]
    sldIdLst.remove(element)
    sldIdLst.insert(int(dst), element)


def main(argv=None):
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    parser = argparse.ArgumentParser(
        description="Edit a .pptx: replace text, update chart data, "
                    "swap images, remove or reorder slides.",
        epilog=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", help="path to the .pptx file")
    parser.add_argument("--output", help="save to this path instead of "
                                         "overwriting the input")
    parser.add_argument("--replace-text", nargs=2, action="append",
                        metavar=("OLD", "NEW"), default=[])
    parser.add_argument("--chart-data", metavar="SPEC_JSON")
    parser.add_argument("--swap-image", nargs=3,
                        metavar=("SLIDE", "SHAPE_NAME", "IMAGE"))
    parser.add_argument("--remove-slide", type=int, metavar="N")
    parser.add_argument("--move-slide", nargs=2, type=int,
                        metavar=("FROM", "TO"))
    args = parser.parse_args(argv)

    prs = Presentation(args.pptx)
    report = {"ok": True, "replacements": 0}

    for old, new in args.replace_text:
        report["replacements"] += replace_text(prs, old, new)
    if args.chart_data:
        update_chart(prs, args.chart_data)
        report["chart_updated"] = True
    if args.swap_image:
        swap_image(prs, *args.swap_image)
        report["image_swapped"] = True
    if args.remove_slide is not None:
        remove_slide(prs, args.remove_slide)
        report["slide_removed"] = args.remove_slide
    if args.move_slide:
        move_slide(prs, *args.move_slide)
        report["slide_moved"] = args.move_slide

    out = args.output or args.pptx
    prs.save(out)
    report["output"] = out
    print(json.dumps(report))
    return 0


if __name__ == "__main__":
    sys.exit(main())
